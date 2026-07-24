import re
from flask import Blueprint, request, jsonify, send_file
from flask_login import login_required
from models import Comment
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image

export_batch_bp = Blueprint('export_batch', __name__, url_prefix='/api/export')

@export_batch_bp.route('/pptx/batch-images', methods=['POST'])
@login_required
def export_batch_with_images():
    """
    Exporta PPTX com Layout Híbrido Automático e Algoritmo Best-Fit:
    - Tabelas Largas (> 1.2 aspect ratio) -> Layout Topo (Full Width)
    - Tabelas Altas/Quadradas -> Layout Lateral (Dashboard)
    - Garante ZERO distorção de imagem (aspect ratio preservado)
    """
    try:
        query_count = int(request.form.get('query_count', 0))
        mes_selecionado = int(request.form.get('mes_selecionado', 0))
        ano_selecionado = int(request.form.get('ano_selecionado', 0))
        periodo_label = f"Resultado {str(mes_selecionado).zfill(2)}/{ano_selecionado}" if mes_selecionado and ano_selecionado else ""

        if query_count == 0:
            return jsonify({'error': 'Nenhuma query fornecida'}), 400
        
        # Criar apresentação
        prs = Presentation()
        prs.slide_width = 12192000  # 13.33 polegadas (16:9 widescreen)
        prs.slide_height = 6858000  # 7.5 polegadas

        # --- FUNÇÃO AUXILIAR: BEST FIT CALCULATOR ---
        # Essa função é o segredo para não distorcer a imagem.
        # Ela calcula o maior tamanho possível mantendo a proporção.
        def get_best_fit_dimensions(img_w, img_h, container_w, container_h, container_left, container_top):
            img_ratio = img_w / img_h
            container_ratio = container_w / container_h
            
            final_w = 0
            final_h = 0

            # Se a imagem é proporcionalmente "mais larga" que o container
            if img_ratio > container_ratio:
                # A largura limita
                final_w = container_w
                final_h = final_w / img_ratio
            else:
                # A altura limita
                final_h = container_h
                final_w = final_h * img_ratio
            
            # Centralizar matematicamente no container
            centered_left = container_left + (container_w - final_w) / 2
            centered_top = container_top + (container_h - final_h) / 2
            
            return centered_left, centered_top, final_w, final_h

        # Threshold para decidir entre Layout A (Topo) e Layout B (Split)
        # Só tabelas claramente panorâmicas (ratio > 1.5) usam Layout A
        THRESHOLD_RATIO = 1.5

        for i in range(query_count):
            try:
                query_id = request.form.get(f'query_id_{i}')
                query_title = request.form.get(f'query_title_{i}')
                table_image_file = request.files.get(f'table_image_{i}')
                is_combo = request.form.get(f'is_combo_{i}', '0') == '1'

                if not all([query_id, query_title, table_image_file]):
                    continue
                
                # Ler imagem para memória
                image_bytes = table_image_file.read()
                image_stream = BytesIO(image_bytes)
                
                # Descobrir dimensões originais com PIL
                with Image.open(image_stream) as img:
                    orig_w, orig_h = img.size
                    img_ratio = orig_w / orig_h
                
                # Resetar ponteiro do stream para o PPTX ler
                image_stream.seek(0)
                
                # Adicionar Slide em Branco
                slide = prs.slides.add_slide(prs.slide_layouts[6])

                # Sem isso o fundo do slide fica no padrão do tema (não
                # necessariamente branco) — a imagem da tabela já vem com
                # fundo branco só ao redor dela mesma, e a combinação parecia
                # uma "ilha" pequena flutuando num fundo diferente
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

                # --- CAPA ---
                if query_id == 'cover':
                    # Usa Best Fit na Capa também para centralizar sem distorcer
                    left, top, w, h = get_best_fit_dimensions(
                        orig_w, orig_h, 
                        prs.slide_width, prs.slide_height, 
                        0, 0
                    )
                    slide.shapes.add_picture(image_stream, left, top, width=w, height=h)
                    continue

                # =========================================================
                # DECISÃO DE LAYOUT BASEADA NA PROPORÇÃO DA IMAGEM
                # =========================================================
                # Página combo (2+ tabelas empilhadas) sempre usa Layout A —
                # a altura maior de várias tabelas juntas derrubava a proporção
                # e fazia cair no Layout B (painel de comentários lateral, que
                # nem faz sentido pra combo), espremendo a imagem sem motivo
                #
                # Layout B só vale a pena quando existe comentário aprovado
                # pra mostrar no painel — sem isso, o painel fica em branco e
                # a tabela é espremida numa coluna estreita sem necessidade
                has_approved_comments = Comment.query.filter_by(
                    query_id=query_id, status='approved'
                ).first() is not None

                if is_combo or img_ratio > THRESHOLD_RATIO or not has_approved_comments:
                    # -----------------------------------------------------
                    # LAYOUT A: FULL WIDTH (Para Tabelas Grandes/Largas)
                    # -----------------------------------------------------
                    MARGIN_X = Inches(0.25)
                    TITLE_TOP = Inches(0.05)
                    TITLE_HEIGHT = Inches(0.6)

                    # 1. Título Topo
                    title_shape = slide.shapes.add_textbox(MARGIN_X, TITLE_TOP, prs.slide_width - (MARGIN_X*2), TITLE_HEIGHT)
                    tf = title_shape.text_frame
                    tf.text = query_title
                    p = tf.paragraphs[0]
                    p.font.size = Pt(22)
                    p.font.bold = True
                    p.font.name = 'Arial'
                    p.font.color.rgb = RGBColor(30, 58, 95)
                    p.alignment = PP_ALIGN.LEFT
                    if periodo_label:
                        p_sub = tf.add_paragraph()
                        p_sub.text = periodo_label
                        p_sub.font.size = Pt(12)
                        p_sub.font.bold = False
                        p_sub.font.name = 'Arial'
                        p_sub.font.color.rgb = RGBColor(100, 116, 139)
                        p_sub.alignment = PP_ALIGN.LEFT

                    # 2. Área Disponível para a Tabela
                    # Começa abaixo do título — a imagem SEMPRE cabe dentro do
                    # slide (contain-fit). Em apresentação (Slideshow) não tem
                    # como rolar pra ver o que passa do slide, então nunca pode
                    # estourar o limite
                    TABLE_TOP = Inches(0.65)

                    available_w = prs.slide_width - (MARGIN_X * 2)
                    available_h = prs.slide_height - TABLE_TOP - MARGIN_X

                    left, top, width, height = get_best_fit_dimensions(
                        orig_w, orig_h, available_w, available_h, MARGIN_X, TABLE_TOP
                    )
                    # Alinha à esquerda e ao topo (em vez de centralizar) — o
                    # espaço que sobra à direita vira coluna de comentários
                    left = MARGIN_X
                    top = TABLE_TOP

                    slide.shapes.add_picture(image_stream, left, top, width=width, height=height)

                    # 3. Comentários — coluna vertical à direita da tabela,
                    # usando o espaço que sobrou. Se a tabela ocupar a largura
                    # toda (não sobra espaço pra coluna), cai pro rodapé
                    sidebar_left = left + width + Inches(0.3)
                    sidebar_w = prs.slide_width - sidebar_left - MARGIN_X
                    if has_approved_comments and sidebar_w >= Inches(2.0):
                        _add_comments_sidebar(slide, query_id, sidebar_left, TABLE_TOP, sidebar_w, height)
                    elif has_approved_comments:
                        footer_top = prs.slide_height - Inches(0.7)
                        _add_comments_footer(slide, query_id, MARGIN_X, footer_top, prs.slide_width - (MARGIN_X*2))

                else:
                    # -----------------------------------------------------
                    # LAYOUT B: Título topo + Tabela esq. + Comentários dir.
                    # Mesma estrutura do Layout A (só entra aqui quando existe
                    # comentário aprovado E a tabela deixa espaço sobrando)
                    # -----------------------------------------------------
                    MARGIN = Inches(0.25)
                    TITLE_TOP = Inches(0.05)
                    TITLE_H = Inches(0.6)
                    TABLE_TOP = Inches(0.65)

                    # 1. Título topo (largura total)
                    title_shape = slide.shapes.add_textbox(
                        MARGIN, TITLE_TOP, prs.slide_width - MARGIN * 2, TITLE_H
                    )
                    tf = title_shape.text_frame
                    tf.word_wrap = True
                    tf.text = query_title
                    p = tf.paragraphs[0]
                    p.font.size = Pt(20)
                    p.font.bold = True
                    p.font.name = 'Arial'
                    p.font.color.rgb = RGBColor(30, 58, 95)
                    p.alignment = PP_ALIGN.LEFT
                    if periodo_label:
                        p_sub = tf.add_paragraph()
                        p_sub.text = periodo_label
                        p_sub.font.size = Pt(11)
                        p_sub.font.bold = False
                        p_sub.font.name = 'Arial'
                        p_sub.font.color.rgb = RGBColor(100, 116, 139)
                        p_sub.alignment = PP_ALIGN.LEFT

                    # 2. Tabela painel esquerdo (contain-fit, alinhada à
                    # esquerda e ao topo — mesma lógica do Layout A)
                    available_w = prs.slide_width - (MARGIN * 2)
                    available_h = prs.slide_height - TABLE_TOP - MARGIN

                    left, top, width, height = get_best_fit_dimensions(
                        orig_w, orig_h, available_w, available_h, MARGIN, TABLE_TOP
                    )
                    left = MARGIN
                    top = TABLE_TOP
                    slide.shapes.add_picture(image_stream, left, top, width=width, height=height)

                    # 3. Comentários painel direito, no espaço que sobrar
                    sidebar_left = left + width + Inches(0.3)
                    sidebar_w = prs.slide_width - sidebar_left - MARGIN
                    _add_comments_sidebar(slide, query_id, sidebar_left, TABLE_TOP, sidebar_w, height)

            except Exception as e:
                print(f"Erro processando query {i}: {e}")
                continue
        
        # Salvar e Enviar
        pptx_bytes = BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)
        
        return send_file(
            pptx_bytes,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name='export_hybrid.pptx'
        )

    except Exception as e:
        return jsonify({'error': str(e)}), 500

# --- Funções Auxiliares para Comentários ---

def _get_comment_text(c):
    """Retorna edited_content se existir, senão content original."""
    return (c.edited_content or c.content).strip()

# "Rótulo: resto do texto" -> rótulo em negrito. Limitado a 40 caracteres antes
# dos ":" pra não pegar um ":" que apareça no meio de uma frase longa qualquer
_LABEL_PREFIX = re.compile(r'^([^:\n]{1,40}):\s*(.*)$', re.DOTALL)

def _add_comment_text_runs(paragraph, text, color, font_size=Pt(11)):
    """Adiciona o texto do comentário à paragraph, deixando em negrito
    qualquer 'Rótulo:' no início do texto (ex: 'FOPAG: texto...')."""
    match = _LABEL_PREFIX.match(text)
    if match:
        label, rest = match.group(1), match.group(2)
        run_label = paragraph.add_run()
        run_label.text = f"{label}: "
        run_label.font.bold = True
        run_label.font.size = font_size
        run_label.font.color.rgb = color

        run_rest = paragraph.add_run()
        run_rest.text = rest
        run_rest.font.bold = False
        run_rest.font.size = font_size
        run_rest.font.color.rgb = color
    else:
        run = paragraph.add_run()
        run.text = text
        run.font.bold = False
        run.font.size = font_size
        run.font.color.rgb = color

def _add_comments_footer(slide, query_id, left, top, width):
    """Adiciona comentários empilhados no rodapé (Layout A) — um por linha."""
    comments = Comment.query.filter_by(query_id=query_id).order_by(Comment.created_at).all()
    approved = [c for c in comments if c.status == 'approved']

    if not approved:
        return

    row_height = Inches(0.28)
    total_height = row_height * len(approved)
    shape = slide.shapes.add_textbox(left, top, width, total_height)
    tf = shape.text_frame
    tf.word_wrap = True

    MAX_CHARS_FOOTER = 160
    for i, c in enumerate(approved):
        author = c.author.username if c.author else 'User'
        text = _get_comment_text(c)
        if len(text) > MAX_CHARS_FOOTER:
            text = text[:MAX_CHARS_FOOTER].rstrip() + '...'
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)

        run_author = p.add_run()
        run_author.text = f"{author}: "
        run_author.font.bold = True
        run_author.font.size = Pt(11)
        run_author.font.color.rgb = RGBColor(30, 58, 95)

        _add_comment_text_runs(p, text, RGBColor(80, 80, 80))

def _add_comments_sidebar(slide, query_id, left, top, width, height):
    """Adiciona comentários numa coluna lateral (Layout A/B) com truncamento controlado."""
    MAX_CHARS = 140
    MAX_VISIBLE = 3

    comments = Comment.query.filter_by(query_id=query_id).order_by(Comment.created_at).all()
    approved = [c for c in comments if c.status == 'approved']

    if not approved:
        return

    shape = slide.shapes.add_textbox(left, top, width, height)
    # Fundo azul bem claro pra dar contraste com o resto do slide (branco)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(235, 243, 252)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)

    p_head = tf.paragraphs[0]
    p_head.text = "Comentários:"
    p_head.font.bold = True
    p_head.font.size = Pt(11)
    p_head.font.color.rgb = RGBColor(80, 80, 80)
    p_head.space_after = Pt(4)

    for c in approved[:MAX_VISIBLE]:
        author = c.author.username if c.author else 'User'
        text = _get_comment_text(c)
        if len(text) > MAX_CHARS:
            text = text[:MAX_CHARS].rstrip() + '...'

        p_auth = tf.add_paragraph()
        p_auth.text = f"{author}:"
        p_auth.font.bold = True
        p_auth.font.size = Pt(11)
        p_auth.font.color.rgb = RGBColor(30, 58, 95)

        p_content = tf.add_paragraph()
        _add_comment_text_runs(p_content, text, RGBColor(60, 60, 60))
        p_content.space_after = Pt(6)

    remaining = len(approved) - MAX_VISIBLE
    if remaining > 0:
        p_more = tf.add_paragraph()
        p_more.text = f"+ {remaining} comentário(s) não exibido(s)"
        p_more.font.size = Pt(11)
        p_more.font.italic = True
        p_more.font.color.rgb = RGBColor(130, 130, 130)