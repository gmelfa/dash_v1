from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from io import BytesIO
from PIL import Image

def create_final_pptx(query_title, table_data, columns, approved_comments):
    """
    Gera PowerPoint para apresentação final
    Contém: Título + Tabela nativa + Comentários aprovados editados
    """
    prs = Presentation()
    # Formato 16:9 widescreen
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Adicionar título
    title_box = slide.shapes.add_textbox(Inches(0.25), Inches(0.3), Inches(9.5), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = query_title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(24)
    title_para.font.bold = True
    
    # Linha separadora
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.25), Inches(0.85),
        Inches(9.5), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line.line.fill.background()
    
    # Criar tabela nativa
    rows_count = len(table_data) + 2  # +2 para cabeçalho duplo
    cols_count = len(columns)
    
    table = slide.shapes.add_table(
        rows_count, cols_count,
        Inches(0.3), Inches(1),
        Inches(9.4), Inches(3.5)
    ).table
    
    # Formatar tabela
    format_table(table, columns, table_data)
    
    # Adicionar comentários aprovados
    if approved_comments:
        comment_top = Inches(4.8)
        add_comments_section(slide, approved_comments, comment_top)
    
    return prs


def format_table(table, columns, data):
    """
    Formata tabela nativa do PowerPoint
    """
    # Primeira linha - Cabeçalho de períodos
    table.cell(0, 0).text = "Em R$ mil"
    table.cell(0, 0).fill.solid()
    table.cell(0, 0).fill.fore_color.rgb = RGBColor(241, 245, 249)
    
    # Merge para "Vertical"
    vertical_cell = table.cell(1, 0)
    vertical_cell.text = "Vertical"
    vertical_cell.fill.solid()
    vertical_cell.fill.fore_color.rgb = RGBColor(31, 56, 100)
    vertical_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    vertical_cell.text_frame.paragraphs[0].font.bold = True
    vertical_cell.text_frame.paragraphs[0].font.size = Pt(9)
    
    # Períodos (10M24 R, 10M25 F, 10M25 R)
    periods = [
        {"name": "10M24 R", "start_col": 1},
        {"name": "10M25 F", "start_col": 4},
        {"name": "10M25 R", "start_col": 7}
    ]
    
    for period in periods:
        cell = table.cell(0, period["start_col"])
        cell.text = period["name"]
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Merge 3 colunas
        if period["start_col"] + 2 < len(columns):
            cell.merge(table.cell(0, period["start_col"] + 2))
    
    # Segunda linha - Nomes das colunas
    for col_idx, col_name in enumerate(columns[1:], 1):  # Pula "Vertical"
        # Remover período entre parênteses
        import re
        clean_name = re.sub(r'\s*\([^)]*\)', '', col_name).strip()
        
        cell = table.cell(1, col_idx)
        cell.text = clean_name
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(31, 56, 100)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(8)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Preencher dados
    for row_idx, row_data in enumerate(data):
        actual_row = row_idx + 2
        
        for col_idx, (col_name, value) in enumerate(zip(columns, row_data.values())):
            cell = table.cell(actual_row, col_idx)
            
            # Formatar valor
            if isinstance(value, (int, float)) and value is not None:
                cell.text = f"{int(value):,}".replace(',', '.')
            else:
                cell.text = str(value) if value is not None else ""
            
            # Detectar subtotais
            vertical_value = list(row_data.values())[0] if row_data else ""
            is_subtotal = any(keyword in str(vertical_value) for keyword in [
                'Operações Premium', 'Total Estratégico', 'Total Op. Premium'
            ])
            
            if is_subtotal:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(180, 198, 231)
                cell.text_frame.paragraphs[0].font.bold = True
            
            cell.text_frame.paragraphs[0].font.size = Pt(7)
            
            # Alinhamento
            if col_idx == 0:
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            else:
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


def add_comments_section(slide, comments, top_position):
    """
    Adiciona seção de comentários no slide
    """
    current_top = top_position
    
    for idx, comment in enumerate(comments, 1):
        # Texto do comentário com autor em negrito
        comment_box = slide.shapes.add_textbox(
            Inches(0.5), current_top,
            Inches(9.0), Inches(0.35)
        )
        comment_box.text_frame.word_wrap = True

        content = comment.get('edited_content') or comment['content']
        author = comment.get('username', 'User')

        p = comment_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT

        run_author = p.add_run()
        run_author.text = f"{author}: "
        run_author.font.bold = True
        run_author.font.size = Pt(11)
        run_author.font.color.rgb = RGBColor(30, 58, 95)

        run_text = p.add_run()
        run_text.text = content
        run_text.font.bold = False
        run_text.font.size = Pt(11)
        run_text.font.color.rgb = RGBColor(60, 60, 60)
        
        current_top += Inches(0.35)

