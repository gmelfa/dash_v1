"""
Teste simples do endpoint batch sem executar queries do Databricks
"""
from flask import Flask, jsonify, request, send_file
from flask_cors import CORS
from flask_login import LoginManager, login_required, current_user
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from io import BytesIO
import json

app = Flask(__name__)
app.config['SECRET_KEY'] = 'test-secret-key'

CORS(app, origins=['http://localhost:5173'], supports_credentials=True)

login_manager = LoginManager()
login_manager.init_app(app)

# Mock user
class MockUser:
    def __init__(self):
        self.id = 1
        self.is_authenticated = True
        
    def get_id(self):
        return str(self.id)

@login_manager.user_loader
def load_user(user_id):
    return MockUser()

@login_manager.request_loader
def load_user_from_request(request):
    return MockUser()

@app.route('/api/test/batch', methods=['POST'])
def test_batch():
    """Teste simples de geração de PPT"""
    try:
        print("DEBUG: Iniciando teste de batch")
        
        # Criar apresentação
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        print("DEBUG: Apresentação criada")
        
        # Adicionar slide de capa
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        print("DEBUG: Slide de capa adicionado")
        
        # Estilo de capa
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(30, 58, 95)
        
        print("DEBUG: Background configurado")
        
        # Adicionar texto
        title_box = slide.shapes.add_textbox(
            Inches(1.1), Inches(3.3), prs.slide_width - Inches(2.2), Inches(1.1)
        )
        text_frame = title_box.text_frame
        text_frame.text = "Teste - Business Review"
        text_frame.paragraphs[0].font.size = Pt(54)
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        text_frame.paragraphs[0].font.bold = True
        
        print("DEBUG: Texto da capa adicionado")
        
        # Salvar
        pptx_bytes = BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)
        
        print(f"DEBUG: PowerPoint salvo, tamanho: {pptx_bytes.getbuffer().nbytes} bytes")
        
        return send_file(
            pptx_bytes,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name='test_batch.pptx'
        )
        
    except Exception as e:
        print(f"ERROR: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    print("Teste do endpoint batch - rodando em http://localhost:5001")
    app.run(debug=True, port=5001)
