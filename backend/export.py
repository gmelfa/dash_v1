from flask import Blueprint, request, jsonify, send_file, current_app
from flask_login import login_required, current_user
from models import Comment
from pptx_service import create_development_pptx, create_final_pptx
import json
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

export_bp = Blueprint('export', __name__, url_prefix='/api/export')

@export_bp.route('/pptx/development', methods=['POST'])
@login_required
def export_development():
    """
    Exporta PowerPoint de desenvolvimento
    Recebe: query_id, query_title, table_image (bytes)
    """
    try:
        query_id = request.form.get('query_id')
        query_title = request.form.get('query_title')
        table_image = request.files.get('table_image')
        
        if not query_id or not query_title or not table_image:
            return jsonify({'error': 'query_id, query_title e table_image são obrigatórios'}), 400
        
        # Ler bytes da imagem
        table_image_bytes = table_image.read()
        
        # Buscar todos os comentários da query
        comments = Comment.query.filter_by(query_id=query_id).order_by(Comment.created_at).all()
        comments_data = [comment.to_dict() for comment in comments]
        
        # Gerar PowerPoint
        prs = create_development_pptx(query_title, table_image_bytes, comments_data)
        
        # Salvar em BytesIO
        pptx_bytes = BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)
        
        # Retornar arquivo
        return send_file(
            pptx_bytes,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=f'{query_id}_development.pptx'
        )
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@export_bp.route('/pptx/final', methods=['POST'])
@login_required
def export_final():
    """
    Exporta PowerPoint final (apenas comentários aprovados com tabela nativa)
    Recebe: query_id, query_title, columns (JSON), table_data (JSON)
    """
    try:
        query_id = request.form.get('query_id')
        query_title = request.form.get('query_title')
        columns_json = request.form.get('columns')
        table_data_json = request.form.get('table_data')
        
        if not all([query_id, query_title, columns_json, table_data_json]):
            return jsonify({'error': 'query_id, query_title, columns e table_data são obrigatórios'}), 400
        
        # Parse JSON
        columns = json.loads(columns_json)
        table_data = json.loads(table_data_json)
        
        # Buscar apenas comentários aprovados
        comments = Comment.query.filter_by(
            query_id=query_id,
            status='approved'
        ).order_by(Comment.created_at).all()
        
        approved_comments = [comment.to_dict() for comment in comments]
        
        # Gerar PowerPoint final com tabela nativa
        prs = create_final_pptx(query_title, table_data, columns, approved_comments)
        
        # Salvar em BytesIO
        pptx_bytes = BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)
        
        # Retornar arquivo
        return send_file(
            pptx_bytes,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=f'{query_id}_final.pptx'
        )
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@export_bp.route('/pptx/batch', methods=['POST'])
@login_required
def export_batch():
    """
    Exporta múltiplas queries em um único PowerPoint
    Recebe: query_ids (array de string)
    Retorna: PowerPoint com um slide por query (cover + queries)
    """
    try:
        data = request.get_json()
        query_ids = data.get('query_ids', [])
        print(f"DEBUG: export_batch called with query_ids: {query_ids}")
        
        if not query_ids or not isinstance(query_ids, list):
            return jsonify({'error': 'query_ids deve ser uma lista de strings'}), 400
        
        # Criar apresentação com dimensões em EMU (inteiros)
        prs = Presentation()
        prs.slide_width = 9144000   # 10 polegadas em EMU (10 * 914400)
        prs.slide_height = 6858000  # 7.5 polegadas em EMU (7.5 * 914400)
        
        # Adicionar slide de capa
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # Estilo de capa
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(30, 58, 95)
        
        # Texto da capa
        title_box = slide.shapes.add_textbox(
            int(Inches(1.1)), int(Inches(3.3)), 
            int(Inches(7.8)), int(Inches(1.1))
        )
        text_frame = title_box.text_frame
        text_frame.text = "Resumo - Grupo SEB (YTD Outubro)"
        text_frame.paragraphs[0].font.size = Pt(54)
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        text_frame.paragraphs[0].font.bold = True

        # Propriedades do arquivo
        prs.core_properties.title = "Resumo - Grupo SEB (YTD Outubro)"
        prs.core_properties.subject = "Resumo - Grupo SEB (YTD Outubro)"
        
        # Reusa o QueryLoader já carregado pelo app.py em vez de criar um novo
        try:
            from app import query_loader
            query_loader.check_for_updates()
            all_queries_data = query_loader.list_queries()
            print(f"DEBUG: Loaded {len(all_queries_data)} queries from QueryLoader")
        except Exception as e:
            print(f"DEBUG: Error loading queries: {e}")
            return jsonify({'error': f'Erro ao carregar queries: {str(e)}'}), 500
        
        # Processar cada query
        for query_id in query_ids:
            try:
                # Encontrar query nos dados carregados
                query_data = next((q for q in all_queries_data if q['id'] == query_id), None)
                if not query_data:
                    print(f"WARN: Query {query_id} not found, skipping")
                    continue
                
                # Executar query no Databricks
                try:
                    print(f"DEBUG: Executing query {query_id}")
                    from databricks import sql
                    import os
                    from dotenv import load_dotenv
                    load_dotenv()
                    
                    conn = sql.connect(
                        server_hostname=os.getenv('DATABRICKS_SERVER_HOSTNAME'),
                        http_path=os.getenv('DATABRICKS_HTTP_PATH'),
                        access_token=os.getenv('DATABRICKS_TOKEN')
                    )
                    cursor = conn.cursor()
                    cursor.execute(query_data['sql_content'])
                    result = cursor.fetchall()
                    columns = [desc[0] for desc in cursor.description]
                    cursor.close()
                    conn.close()
                    print(f"DEBUG: Query {query_id} executed, rows: {len(result)}")
                    
                    # Buscar comentários
                    comments = Comment.query.filter_by(query_id=query_id).order_by(Comment.created_at).all()
                    comments_data = [comment.to_dict() for comment in comments]
                    
                    # Criar slide
                    slide_layout = prs.slide_layouts[5]
                    new_slide = prs.slides.add_slide(slide_layout)
                    
                    # Título do slide
                    title_shape = new_slide.shapes.add_textbox(
                        int(Inches(1)), int(Inches(0.5)), 
                        int(Inches(8)), int(Inches(0.5))
                    )
                    title_frame = title_shape.text_frame
                    title_frame.text = query_data['name']
                    title_frame.paragraphs[0].font.size = Pt(32)
                    title_frame.paragraphs[0].font.bold = True
                    title_frame.paragraphs[0].font.color.rgb = RGBColor(30, 58, 95)
                    
                    # Tabela (limitar a 14 linhas de dados)
                    rows = min(len(result) + 1, 15)
                    cols = len(columns)
                    table_shape = new_slide.shapes.add_table(
                        rows, cols,
                        int(Inches(1)), int(Inches(1.3)),
                        int(Inches(8)), int(Inches(4.9))
                    ).table
                    
                    # Cabeçalho da tabela
                    for col_idx, column_name in enumerate(columns):
                        cell = table_shape.cell(0, col_idx)
                        cell.text = str(column_name)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(30, 58, 95)
                        paragraph = cell.text_frame.paragraphs[0]
                        paragraph.font.color.rgb = RGBColor(255, 255, 255)
                        paragraph.font.bold = True
                        paragraph.font.size = Pt(11)
                    
                    # Dados da tabela
                    for row_idx, row_data in enumerate(result[:14], 1):
                        for col_idx, value in enumerate(row_data):
                            cell = table_shape.cell(row_idx, col_idx)
                            cell.text = str(value) if value is not None else ""
                            if row_idx % 2 == 0:
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = RGBColor(248, 250, 252)
                            paragraph = cell.text_frame.paragraphs[0]
                            paragraph.font.size = Pt(10)
                    
                    # Informação de comentários
                    if comments_data:
                        comment_shape = new_slide.shapes.add_textbox(
                            int(Inches(1)), int(Inches(6.45)), 
                            int(Inches(8)), int(Inches(0.55))
                        )
                        comment_frame = comment_shape.text_frame
                        comment_frame.word_wrap = True
                        pending_count = sum(1 for c in comments_data if c.get('status') == 'pending')
                        comment_frame.text = f"📝 {len(comments_data)} comentário(s) ({pending_count} pendente(s))"
                        comment_frame.paragraphs[0].font.size = Pt(9)
                        comment_frame.paragraphs[0].font.italic = True
                        comment_frame.paragraphs[0].font.color.rgb = RGBColor(100, 116, 139)
                    
                except Exception as e:
                    print(f"ERROR executing query {query_id}: {str(e)}")
                    import traceback
                    traceback.print_exc()
                    continue
                    
            except Exception as e:
                print(f"ERROR processing query {query_id}: {str(e)}")
                continue
        
        # Salvar PowerPoint
        pptx_bytes = BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)
        
        print(f"DEBUG: PowerPoint generated, size: {pptx_bytes.getbuffer().nbytes} bytes")
        
        return send_file(
            pptx_bytes,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name='export_batch.pptx'
        )
        
    except Exception as e:
        print(f"ERROR in export_batch: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500

